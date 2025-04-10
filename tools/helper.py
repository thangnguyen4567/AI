from datetime import datetime
import pytz
import logging
import os
from queue import Queue
from threading import Thread
import asyncio
import importlib
import random
import string
from pathlib import Path
from langchain_community.document_loaders import PyPDFLoader
import requests
from langchain_community.document_loaders import Docx2txtLoader
from pptx import Presentation

def convert_unixtime(unixtime):
    unixtime = int(unixtime)
    utc_datetime = datetime.utcfromtimestamp(unixtime)
    utc_timezone = pytz.timezone('UTC')
    utc_datetime = utc_timezone.localize(utc_datetime)
    vn_timezone = pytz.timezone('Asia/Ho_Chi_Minh')
    return utc_datetime.astimezone(vn_timezone).strftime("%d/%m/%Y %H:%M")

def logging_data(info):
    log_file_path = os.path.join('logs', 'get_table.log')
    logging.basicConfig(filename=log_file_path, level=logging.INFO, format='%(asctime)s - %(levelname)s \n %(message)s')
    string_log = ''
    for value in info:
        string_log +=value + '\n'
    logging.info(string_log)
        
def remove_stopwords(text):
    file_path = './static/vietnamese-stopwords.txt'
    with open(file_path, 'r', encoding='utf-8') as file:
        stopwords = set(line.strip().lower() for line in file)

    words = text.split()  # Tách từ trong văn bản
    filtered_words = [word for word in words if word.lower() not in stopwords]
    filtered_text = ' '.join(filtered_words)  # Ghép lại các từ đã lọc
    return filtered_text

def async_to_sync(generator_func):
    def sync_generator(*args, **kwargs):
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        gen = generator_func(*args, **kwargs)
        q = Queue()

        def enqueue_data():
            not_done = True
            while not_done:
                try:
                    item = loop.run_until_complete(gen.__anext__())
                    q.put(item)
                except StopAsyncIteration:
                    not_done = False
            q.put(StopIteration)

        thread = Thread(target=enqueue_data)
        thread.start()

        while True:
            item = q.get()
            if item == StopIteration:
                break
            if isinstance(item, list):
                metadata = item[0]
            else:
                yield item

        loop.call_soon_threadsafe(loop.stop)
        thread.join()

    return sync_generator

def get_chatbot(data):
    try:
        module_name = f"project.{data['context']}.services.chatbot"
        chatbot_module = importlib.import_module(module_name)
        return chatbot_module.ChatBot(data)
    except Exception as e:
        print(e)
    
def generate_random_string(length=10):
    letters = string.ascii_letters
    return ''.join(random.choice(letters) for i in range(length))


def file_reader(file_path):

    path = Path(file_path)
    typefile = path.suffix.lower()
    text = ''

    if typefile == '.pdf':
        docs = PyPDFLoader(file_path)
        for doc in docs.load():
            text += doc.page_content

    elif typefile == '.docx':
        response = requests.get(file_path)
        random_string = generate_random_string()
        name = random_string+'.docx'
        with open(name, 'wb') as file:
            file.write(response.content)
        docs = Docx2txtLoader(name)
        text = docs.load()[0].page_content
        os.remove(name)

    elif typefile == '.pptx':
        response = requests.get(file_path)
        random_string = generate_random_string()
        name = random_string+'.pptx'
        with open(name, 'wb') as file:
            file.write(response.content)
            
        presentation = Presentation(name)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
        os.remove(name)

    return text