from factory.base.training import Training
from langchain_community.document_loaders import PyPDFLoader
from langchain_community.document_loaders import UnstructuredURLLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from langchain_community.document_loaders import Docx2txtLoader
from pptx import Presentation
from pathlib import Path
import requests
import tempfile
import os
from tools.helper import generate_random_string

class TrainingResourceSaas(Training):
    def __init__(self):
        super().__init__()
        self.text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
        self.columns = ['title','content','coursemoduleid','courseid','source','coursename']
        self.redis_client = self.vector_db.connect_client()

    def save_training_data(self,data):

        finaldocx = []
        path = Path(data['source'])
        typefile = path.suffix.lower()
        collection = self.get_collection_name(data)

        cookies = {}
        cookies['tenantid'] = data['contextdata']['collection']
        response = requests.get(data['source'], cookies=cookies)

        with tempfile.NamedTemporaryFile(delete=False) as temp_file:
            temp_file.write(response.content)
            temp_path = temp_file.name
        
        try:
            
            if typefile == '.pdf':
                docs = PyPDFLoader(temp_path)
                text = ''
                for doc in docs.load():
                    text += doc.page_content
                all_splits = [Document(page_content=split, metadata={}) for split in self.text_splitter.split_text(text)]

            elif typefile == '.docx':

                docs = Docx2txtLoader(temp_path)
                all_splits = self.text_splitter.split_documents(docs.load())

            elif typefile == '.pptx':

                full_text = ''
                presentation = Presentation(temp_path)
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            full_text += shape.text

                all_splits = self.text_splitter.split_text(full_text)

            else:
                docs = UnstructuredURLLoader(urls=[data['source']])
                all_splits = self.text_splitter.split_documents(docs.load())

            metadata = {}

            for key,value in data.items():
                if key in self.columns:
                    metadata[key] = value
                    
            for key in self.redis_client.scan_iter("doc:"+collection+"*"):
                coursemoduleid = self.redis_client.hget(key,'coursemoduleid').decode()
                if coursemoduleid == metadata['coursemoduleid']:
                    self.redis_client.delete(key)
                    
            for doc in all_splits:
                content = 'Tài liệu: ' + metadata['title']
                content += 'Thuộc lớp học: ' + metadata['coursename'] 
                if hasattr(doc, 'page_content'):
                    content += doc.page_content
                else:
                    content += doc
                finaldocx.append(Document(page_content=content,metadata=metadata))

            os.remove(temp_path)
            self.vector_db.add_vectordb(finaldocx,collection)

            self.reponse['error'] = False
            self.reponse['message'] = 'Training thành công'

            return self.reponse
        
        except Exception as e:

            print(str(e))

            self.reponse['error'] = True
            self.reponse['message'] = f'Training thất bại: {str(e)}'

            return self.reponse

    def delete_training_data(self,data):

        collection = self.get_collection_name(data)

        try:
            coursemoduleids_set = set(data['coursemoduleids'])
            keys_to_delete = []

            for key in self.redis_client.scan_iter("doc:"+collection+"*"):
                coursemoduleid_redis = self.redis_client.hget(key, 'coursemoduleid').decode()
                if coursemoduleid_redis in coursemoduleids_set:
                    keys_to_delete.append(key)

            if keys_to_delete:
                self.redis_client.delete(*keys_to_delete)

                self.reponse['error'] = False
                self.reponse['message'] = 'Xóa dữ liệu thành công'

            return self.reponse
                        
        except Exception as e:

            print(str(e))
            self.reponse['error'] = True
            self.reponse['message'] = f'Xóa dữ liệu thất bại: {str(e)}'

            return self.reponse
        
    def check_training_duplication(self):
        pass

    def get_collection_name(self,data):
        return 'resource_'+data['contextdata']['collection']
